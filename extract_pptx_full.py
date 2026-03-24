from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os

prs = Presentation(r'C:\Users\tamirdresher\source\repos\tamresearch1\Final_project_Presentation_real.pptx')

print(f'Total slides: {len(prs.slides)}')
print(f'Slide size: {prs.slide_width.inches:.1f}" x {prs.slide_height.inches:.1f}"')
print()

# Also extract images to inspect visually
img_dir = r'C:\Users\tamirdresher\source\repos\tamresearch1\pptx_images'
os.makedirs(img_dir, exist_ok=True)

img_count = 0
for i, slide in enumerate(prs.slides):
    print(f'=== SLIDE {i+1} (layout: {slide.slide_layout.name}) ===')
    print(f'  Shapes count: {len(slide.shapes)}')
    for shape in slide.shapes:
        print(f'  Shape: {shape.name} | type={shape.shape_type} | pos=({shape.left/914400:.2f}", {shape.top/914400:.2f}") | size=({shape.width/914400:.2f}"x{shape.height/914400:.2f}")')
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    print(f'    TEXT: {text}')
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            img_count += 1
            try:
                image = shape.image
                ext = image.ext
                img_path = os.path.join(img_dir, f'slide{i+1}_img{img_count}.{ext}')
                with open(img_path, 'wb') as f:
                    f.write(image.blob)
                print(f'    IMAGE saved: slide{i+1}_img{img_count}.{ext} ({len(image.blob)} bytes)')
            except Exception as e:
                print(f'    IMAGE error: {e}')
        if hasattr(shape, 'table'):
            print(f'    TABLE:')
            for row in shape.table.rows:
                row_data = [cell.text.strip() for cell in row.cells]
                print('      ' + ' | '.join(row_data))
        # Check for group shapes
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            print(f'    GROUP shape with {len(shape.shapes)} sub-shapes')
            for sub in shape.shapes:
                if sub.has_text_frame:
                    for para in sub.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            print(f'      GROUP TEXT: {text}')
    print()

print(f'\nTotal images extracted: {img_count}')
