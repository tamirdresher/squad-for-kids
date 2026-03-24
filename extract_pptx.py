from pptx import Presentation

prs = Presentation(r'C:\Users\tamirdresher\source\repos\tamresearch1\Final_project_Presentation_real.pptx')

print(f'Total slides: {len(prs.slides)}')
print(f'Slide size: {prs.slide_width.inches:.1f}" x {prs.slide_height.inches:.1f}"')
print()

for i, slide in enumerate(prs.slides):
    print(f'=== SLIDE {i+1} (layout: {slide.slide_layout.name}) ===')
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    print(f'  [{shape.name}]: {text}')
        if hasattr(shape, 'table'):
            print(f'  [TABLE {shape.name}]:')
            for row in shape.table.rows:
                row_data = [cell.text.strip() for cell in row.cells]
                print('    ' + ' | '.join(row_data))
    print()
